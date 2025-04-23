import urllib.request
import csv
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from collections import defaultdict
from PIL import Image

def limpar_nome_arquivo(nome):
    return re.sub(r'[\\/*?:"<>|\n\r]', "_", nome)

def selecionar_arquivo():
    caminho = filedialog.askopenfilename(
        title="Selecione o arquivo CSV",
        filetypes=(("Arquivos CSV", "*.csv"), ("Todos os arquivos", "*.*"))
    )
    if caminho:
        entry_arquivo.delete(0, tk.END)
        entry_arquivo.insert(0, caminho)
        carregar_colunas(caminho)

def carregar_colunas(caminho):
    try:
        with open(caminho, 'r', newline='', encoding='utf-8') as file:
            reader = csv.DictReader(file, delimiter=';')
            colunas = reader.fieldnames

            combo_pessoa['values'] = colunas
            combo_data['values'] = colunas
            combo_local['values'] = colunas
            for combo in [combo_imagem1, combo_imagem2, combo_imagem3]:
                combo['values'] = colunas
            listbox_info.delete(0, tk.END)
            for col in colunas:
                listbox_info.insert(tk.END, col)
    except Exception as e:
        messagebox.showerror("Erro", f"Não foi possível ler o arquivo:\n{e}")

def gerar_apresentacao():
    caminho_csv = entry_arquivo.get()
    nome_arquivo_ppt = entry_nome_ppt.get().strip() or "apresentacao_fotos"
    nome_arquivo_ppt = nome_arquivo_ppt.replace(" ", "_") + ".pptx"

    col_pessoa = combo_pessoa.get()
    col_data = combo_data.get()
    col_local = combo_local.get()
    col_imagens = [combo.get() for combo in [combo_imagem1, combo_imagem2, combo_imagem3] if combo.get()]
    col_infos = [listbox_info.get(i) for i in listbox_info.curselection()]

    try:
        fotos_por_slide = min(4, int(entry_fotos_por_slide.get()))
    except:
        messagebox.showwarning("Aviso", "Número de fotos por slide inválido.")
        return

    if not caminho_csv or not col_pessoa or not col_data or not col_local or not col_imagens:
        messagebox.showwarning("Atenção", "Preencha todos os campos obrigatórios.")
        return

    if not os.path.exists('fotos'):
        os.makedirs('fotos')

    prs = Presentation()
    grupos = defaultdict(list)

    with open(caminho_csv, 'r', newline='', encoding='utf-8') as file:
        reader = csv.DictReader(file, delimiter=';')
        for row in reader:
            pessoa = row[col_pessoa]
            data = row[col_data]
            local = row[col_local]
            grupo_id = f"{pessoa} - {data}\nLocal: {local}"
            imagens_info = []
            for col_img in col_imagens:
                if row[col_img]:
                    imagens_info.append((row[col_img], {col: row[col] for col in col_infos}))
            if imagens_info:
                grupos[grupo_id].extend(imagens_info)

    for grupo, imagens in grupos.items():
        for i in range(0, len(imagens), fotos_por_slide):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
            tf_titulo = titulo.text_frame
            tf_titulo.text = grupo
            tf_titulo.paragraphs[0].font.size = Pt(20)

            imagens_do_grupo = imagens[i:i+fotos_por_slide]
            total = len(imagens_do_grupo)
            espacamento = Inches(0.3)
            altura_max = Inches(4.5)
            y = Inches(1.0)

            larguras = []
            imagens_temp = []
            for idx, (url, _) in enumerate(imagens_do_grupo):
                nome_arquivo_seguro = limpar_nome_arquivo(f"img_{grupo}_{i+idx}.jpg")
                nome_img = os.path.join("fotos", nome_arquivo_seguro)
                urllib.request.urlretrieve(url, nome_img)
                img = Image.open(nome_img)
                ratio = altura_max / img.height
                nova_largura = img.width * ratio
                larguras.append(nova_largura)
                imagens_temp.append((nome_img, nova_largura))
                img.close()

            largura_total = sum(larguras) + espacamento * (total - 1)
            x_start = (Inches(10) - largura_total) / 2
            x = x_start

            for idx, ((url, info), (nome_img, largura)) in enumerate(zip(imagens_do_grupo, imagens_temp)):
                try:
                    slide.shapes.add_picture(nome_img, x, y, height=altura_max)

                    caixa = slide.shapes.add_textbox(x, y + altura_max + Inches(0.1), largura, Inches(0.6))
                    tf = caixa.text_frame
                    tf.clear()
                    for col in col_infos:
                        p = tf.add_paragraph()
                        p.text = f"{col}: {info.get(col, '')}"
                        p.font.size = Pt(10)

                    x += largura + espacamento
                except Exception as e:
                    print(f"Erro ao baixar/adicionar imagem: {e}")

    prs.save(nome_arquivo_ppt)
    messagebox.showinfo("Sucesso", f"Apresentação '{nome_arquivo_ppt}' gerada com sucesso!")

# Interface tkinter
janela = tk.Tk()
janela.title("Gerador de PPT - InStore")
janela.geometry("650x660")
janela.resizable(False, False)

tk.Label(janela, text="Arquivo CSV:").pack(pady=(10,0))
frame = tk.Frame(janela)
frame.pack()
entry_arquivo = tk.Entry(frame, width=50)
entry_arquivo.pack(side=tk.LEFT, padx=5)
tk.Button(frame, text="Selecionar", command=selecionar_arquivo).pack(side=tk.LEFT)

tk.Label(janela, text="Coluna: Pessoa").pack()
combo_pessoa = ttk.Combobox(janela, state="readonly", width=50)
combo_pessoa.pack()

tk.Label(janela, text="Coluna: Data").pack()
combo_data = ttk.Combobox(janela, state="readonly", width=50)
combo_data.pack()

tk.Label(janela, text="Coluna: Local").pack()
combo_local = ttk.Combobox(janela, state="readonly", width=50)
combo_local.pack()

tk.Label(janela, text="Colunas de Imagem:").pack()
combo_imagem1 = ttk.Combobox(janela, state="readonly", width=50)
combo_imagem1.pack(pady=(2,0))
combo_imagem2 = ttk.Combobox(janela, state="readonly", width=50)
combo_imagem2.pack(pady=(2,0))
combo_imagem3 = ttk.Combobox(janela, state="readonly", width=50)
combo_imagem3.pack(pady=(2,10))

tk.Label(janela, text="Colunas de Informação:").pack()
listbox_info = tk.Listbox(janela, selectmode=tk.MULTIPLE, width=60, height=5)
listbox_info.pack()

tk.Label(janela, text="Fotos por slide (máx. 4):").pack()
entry_fotos_por_slide = tk.Entry(janela, width=10)
entry_fotos_por_slide.insert(0, "4")
entry_fotos_por_slide.pack()

tk.Label(janela, text="Nome do arquivo (sem .pptx):").pack(pady=(10,0))
entry_nome_ppt = tk.Entry(janela, width=50)
entry_nome_ppt.pack()

tk.Button(janela, text="Gerar PowerPoint", bg="#4CAF50", fg="white", font=("Arial", 12), command=gerar_apresentacao).pack(pady=20)

janela.mainloop()
